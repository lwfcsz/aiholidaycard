from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(26, 26, 46)
    background.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_list):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()
    
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(102, 126, 234)
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(16)

add_title_slide(prs, "AI节日贺卡生成系统", "智能祝福语生成 + 电子贺卡模板")

add_content_slide(prs, "项目背景", [
    "🎯 市场需求分析",
    "   • 节日祝福场景广泛：生日、春节、毕业等",
    "   • 传统祝福语缺乏新意，难以定制化",
    "   • 电子贺卡分享便捷，但内容同质化严重",
    "",
    "💡 项目目标",
    "   • 融合AI技术的电子贺卡生成平台",
    "   • 智能生成多场景、多风格的祝福语",
    "   • 美观的贺卡模板和一键分享功能"
])

add_content_slide(prs, "项目简介", [
    "🎨 核心功能",
    "   • AI智能生成祝福语 | 多场景覆盖 | 一键生成贺卡 | 便捷分享",
    "",
    "✨ 模板风格",
    "   • 温馨风格：温暖感人，适合亲朋好友",
    "   • 幽默风格：轻松搞笑，适合熟人朋友",
    "   • 优雅风格：文雅得体，适合正式场合",
    "",
    "🔗 分享功能",
    "   • 生成唯一分享链接 | 支持复制和系统分享 | 随时随地查看"
])

add_content_slide(prs, "技术栈", [
    "⚙️ 后端技术",
    "   • Python 3.8+ | FastAPI 高性能框架 | RESTful API 设计",
    "",
    "🎭 前端技术",
    "   • Vue 3 响应式框架 | Tailwind CSS 美化 | CDN 方式加载依赖",
    "",
    "🔧 其他技术",
    "   • JSON 数据格式 | CORS 跨域支持 | UUID 唯一标识生成",
    "",
    "🏗️ 系统架构",
    "   • 前端(Vue 3) → API请求 → 后端(FastAPI) → 祝福语生成"
])

add_content_slide(prs, "系统架构", [
    "📱 前端层（表现层）",
    "   • Vue 3 响应式框架 | Tailwind CSS 美化界面",
    "   • 用户交互界面设计 | 贺卡实时预览功能",
    "",
    "⚙️ 后端层（业务逻辑层）",
    "   • FastAPI 高性能框架 | RESTful API 接口设计",
    "   • 祝福语模板库管理 | 贺卡数据存储",
    "",
    "🤖 AI层（智能服务层）",
    "   • 智能祝福语生成 | 多风格模板匹配",
    "   • 场景智能识别 | 个性化内容定制"
])

add_content_slide(prs, "核心功能", [
    "🎂 生日祝福",
    "   • 温馨风格：愿你的每一天都像生日蛋糕一样甜蜜",
    "   • 幽默风格：又老了一岁！颜值还是在线的！",
    "   • 优雅风格：良辰吉日，福寿安康",
    "",
    "🏮 节日祝福",
    "   • 新春祝福：恭喜发财，万事如意",
    "   • 圣诞祝福：银装素裹，圣诞祥和",
    "",
    "🎓 毕业祝福",
    "   • 前程似锦，归来仍是少年",
    "   • 这是终点，也是起点，未来的路勇敢去闯"
])

add_content_slide(prs, "贺卡模板", [
    "💕 温馨风格",
    "   • 渐变色：粉色到红色 | 温暖感人，适合亲友",
    "",
    "😄 幽默风格",
    "   • 渐变色：粉红到金黄 | 轻松活泼，趣味十足",
    "",
    "🎩 优雅风格",
    "   • 渐变色：蓝色到青色 | 文雅得体，清新脱俗",
    "",
    "🎉 庆典风格",
    "   • 渐变色：多彩渐变 | 喜庆热闹，活力四射"
])

add_content_slide(prs, "项目特色", [
    "🤖 AI智能化",
    "   • 内置丰富的祝福语模板库",
    "   • 支持多种风格智能匹配",
    "   • 可扩展接入真实AI大模型",
    "",
    "📱 响应式设计",
    "   • 采用 Vue 3 + Tailwind CSS 构建",
    "   • 完美适配各种设备屏幕",
    "",
    "🔗 一键分享",
    "   • 生成唯一分享链接",
    "   • 支持一键复制和系统分享"
])

add_content_slide(prs, "部署方式", [
    "💻 方式一：本地运行（推荐）",
    "   1. 启动后端：cd backend → pip install -r requirements.txt → python app.py",
    "   2. 访问地址：http://localhost:8000",
    "",
    "🚀 方式二：一键启动",
    "   • 双击运行 一键启动.bat",
    "   • 自动安装依赖并启动服务",
    "",
    "🐳 方式三：Docker部署",
    "   • 编写 docker-compose.yml",
    "   • 一键启动所有服务，容器化便于管理",
    "",
    "☁️ 方式四：云服务器部署",
    "   • Nginx 反向代理 | Gunicorn + Uvicorn | 支持高并发"
])

add_title_slide(prs, "总结", "感谢观看！")

prs.save('e:\\保存\\pythonweb全栈期末作业\\PPT项目讲解.pptx')
print("PPT文件已生成：e:\\保存\\pythonweb全栈期末作业\\PPT项目讲解.pptx")
